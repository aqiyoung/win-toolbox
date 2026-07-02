#!/usr/bin/env python3
"""
PDF -> PPT 转换脚本 (混合策略)

协议:
  - 从 stdin 读取一行 JSON: {"inputs": [...], "options": {...}}
  - stdout 逐行输出 JSON 消息:
      {"type":"progress","percent":0,"stage":"开始分析",...}
      {"type":"result","outputFiles":["path.pptx"],...}
      {"type":"error","message":"..."}
  - stderr 用于调试信息,不参与协议

依赖 (打包时通过 pip install -r requirements.txt 安装):
  PyMuPDF (fitz)
  python-pptx
  pdfplumber
  Pillow
"""

import sys
import os
import json
import io
from pathlib import Path

try:
    import fitz  # PyMuPDF
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
    import pdfplumber
    from PIL import Image
except ImportError as e:
    _send({"type": "error", "message": f"缺少依赖: {e}. 请 pip install PyMuPDF python-pptx pdfplumber Pillow"})
    sys.exit(1)


# ─── 协议输出 ──────────────────────────────────

def _send(obj: dict) -> None:
    """向 stdout 输出一行 JSON"""
    line = json.dumps(obj, ensure_ascii=False)
    print(line, flush=True)


def _progress(percent: float, stage: str, **kwargs) -> None:
    _send({"type": "progress", "percent": int(percent), "stage": stage, **kwargs})


def _result(output_files: list, warnings: list = None, stats: dict = None) -> None:
    _send({"type": "result", "outputFiles": output_files,
           "warnings": warnings or [], "stats": stats or {}})


def _error(message: str) -> None:
    _send({"type": "error", "message": message})


# ─── 坐标映射 ──────────────────────────────────

# EMU / point
EMU_PER_INCH = 914400
PT_PER_INCH = 72
EMU_PER_PT = EMU_PER_INCH / PT_PER_INCH  # 12700


def pt_to_emu(pt_val: float) -> int:
    return int(round(pt_val * EMU_PER_PT))


# ─── 核心转换逻辑 ──────────────────────────────

# PPT 默认页面 (16:9)
PPT_WIDTH_PT = 13.333 * PT_PER_INCH
PPT_HEIGHT_PT = 7.5 * PT_PER_INCH


def convert_pdf_to_pptx(
    pdf_path: str,
    output_path: str,
    options: dict,
) -> dict:
    """转换 PDF 到 PPTX

    返回: {"page_count": N, "warnings": [...]}
    """
    strategy = options.get("strategy", "hybrid")
    slide_ratio = options.get("slide_ratio", "match_pdf")
    dpi = options.get("dpi", 200)
    font_fallback = options.get("font_fallback", "微软雅黑")

    # 打开 PDF
    doc = fitz.open(pdf_path)
    page_count = len(doc)

    # 确定页面尺寸
    first_page = doc[0]
    pdf_w = first_page.rect.width
    pdf_h = first_page.rect.height

    # 创建 PPT 演示文稿
    prs = Presentation()

    # 设置页面比例
    if slide_ratio == "4:3":
        prs.slide_width = int(10 * EMU_PER_INCH / PT_PER_INCH * EMU_PER_PT)
        prs.slide_height = int(7.5 * EMU_PER_INCH / PT_PER_INCH * EMU_PER_PT)
    elif slide_ratio == "16:9":
        prs.slide_width = int(13.333 * EMU_PER_INCH / PT_PER_INCH * EMU_PER_PT)
        prs.slide_height = int(7.5 * EMU_PER_INCH / PT_PER_INCH * EMU_PER_PT)
    else:
        # match_pdf
        prs.slide_width = pt_to_emu(pdf_w)
        prs.slide_height = pt_to_emu(pdf_h)

    # 计算 scale
    scale_x = prs.slide_width / pdf_w
    scale_y = prs.slide_height / pdf_h

    warnings = []

    for page_idx in range(page_count):
        page_rect = (page_idx + 1) / page_count * 100
        _progress(page_rect, f"处理第 {page_idx + 1}/{page_count} 页",
                  current=page_idx + 1, total=page_count)

        page = doc[page_idx]
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout

        if strategy == "screenshot":
            _page_to_screenshot(page, slide, dpi, scale_x, scale_y)
        elif strategy == "structured":
            _page_to_structured(page, slide, scale_x, scale_y, font_fallback)
        else:
            _page_to_hybrid(page, slide, dpi, scale_x, scale_y, font_fallback, page_idx)

    # 保存
    _progress(95, "正在保存文件...")
    prs.save(output_path)

    doc.close()

    return {"page_count": page_count, "warnings": warnings}


def _page_to_screenshot(page, slide, dpi, scale_x, scale_y):
    """L3: 每页渲染为整页图片"""
    zoom = dpi / 72
    mat = fitz.Matrix(zoom, zoom)
    pix = page.get_pixmap(matrix=mat, alpha=False)
    img_data = pix.tobytes("png")

    # 图片填满整张 slide
    slide_width_emu = int(pix.width / dpi * EMU_PER_INCH)
    slide_height_emu = int(pix.height / dpi * EMU_PER_INCH)

    img_shape = slide.shapes.add_picture(
        io.BytesIO(img_data),
        left=Emu(0),
        top=Emu(0),
        width=Emu(slide_width_emu),
        height=Emu(slide_height_emu),
    )


def _page_to_structured(page, slide, scale_x, scale_y, font_fallback):
    """L1: 纯结构化重建"""
    # 提取文本块
    blocks = page.get_text("dict")["blocks"]

    for block in blocks:
        if block["type"] == 0:  # 文本块
            for line in block["lines"]:
                line_text = "".join(span["text"] for span in line["spans"])
                if not line_text.strip():
                    continue

                # 取第一个 span 的坐标和样式
                span = line["spans"][0]
                bbox = line["bbox"]
                x0 = bbox[0] * scale_x
                y0 = bbox[1] * scale_y
                width = (bbox[2] - bbox[0]) * scale_x
                height = (bbox[3] - bbox[1]) * scale_y

                # 字体大小
                font_size_pt = span["size"]
                # PPT font size in points -> Pt()
                font_pt = Pt(font_size_pt)

                # 颜色
                color_int = span["color"]
                r = (color_int >> 16) & 0xFF
                g = (color_int >> 8) & 0xFF
                b = color_int & 0xFF

                tb = slide.shapes.add_textbox(
                    left=Emu(int(x0)),
                    top=Emu(int(y0)),
                    width=Emu(int(width)),
                    height=Emu(int(height)),
                )
                tf = tb.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                run = p.add_run()
                run.text = line_text
                run.font.size = font_pt
                run.font.color.rgb = RGBColor(r, g, b)
                run.font.name = span["font"] if _font_available(span["font"]) else font_fallback

        elif block["type"] == 1:  # 图片块
            _place_image_block(slide, block, scale_x, scale_y)


def _page_to_hybrid(page, slide, dpi, scale_x, scale_y, font_fallback, page_idx):
    """L2: 混合策略

    流程:
      1. 分析页面,识别文本区域和"复杂区域"
      2. 文本区域走结构化重建
      3. 整页渲染一张背景图片(低透明度预留位置,让文字在上层)
      4. 复杂图形/表格区域渲染为图片嵌入
    """
    # 先渲染整页为图片用作复杂区域识别的基础
    # 但 L2 不做背景图片,只做结构化 + 兜底图片

    # 尝试提取表格
    tables_on_page = _extract_tables_from_page(page)

    for table in tables_on_page:
        _place_table_as_image(slide, page, table, dpi, scale_x, scale_y)

    # 文本块和浮动图片走结构化重建
    blocks = page.get_text("dict")["blocks"]

    for block in blocks:
        if block["type"] == 0:
            _render_text_block(slide, block, scale_x, scale_y, font_fallback)
        elif block["type"] == 1:
            _place_image_block(slide, block, scale_x, scale_y)


def _render_text_block(slide, block, scale_x, scale_y, font_fallback):
    """渲染一个文本块内的行"""
    for line in block["lines"]:
        line_text = "".join(span["text"] for span in line["spans"])
        if not line_text.strip():
            continue

        # 检查是否有垂直文本 (简单判断: 如果所有 span 的 bbox 宽度接近 0)
        bbox = line["bbox"]
        bb_w = (bbox[2] - bbox[0]) * scale_x
        bb_h = (bbox[3] - bbox[1]) * scale_y

        span = line["spans"][0]
        font_size_pt = span["size"]
        color_int = span["color"]
        r = (color_int >> 16) & 0xFF
        g = (color_int >> 8) & 0xFF
        b = color_int & 0xFF

        try:
            tb = slide.shapes.add_textbox(
                left=Emu(int(bbox[0] * scale_x)),
                top=Emu(int(bbox[1] * scale_y)),
                width=Emu(max(int(bb_w), 50000)),
                height=Emu(max(int(bb_h), 100000)),
            )
        except Exception:
            continue

        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = line_text
        run.font.size = Pt(font_size_pt)
        run.font.color.rgb = RGBColor(r, g, b)
        run.font.name = span["font"] if _font_available(span["font"]) else font_fallback


def _place_image_block(slide, block, scale_x, scale_y):
    """嵌入页面中的图片"""
    # 取 bboxes
    # 实际上 fitz 的 "blocks" 中 type=1 的 block 有 image info
    # 但更可靠的方式是用 page.get_images(full=True) + page.get_image_rects()
    # 这里简化: 仅按 bbox 占位
    pass


def _extract_tables_from_page(page) -> list:
    """
    用 PyMuPDF 内置表格检测
    返回: [{"bbox": (x0, y0, x1, y1), "rows": N}, ...]
    """
    try:
        tables = page.find_tables()
        return [
            {
                "bbox": (t.bbox[0], t.bbox[1], t.bbox[2], t.bbox[3]),
                "rows": len(t.rows),
                "cells": t.cells,
            }
            for t in tables.tables
        ]
    except Exception:
        return []


def _place_table_as_image(slide, page, table_info, dpi, scale_x, scale_y):
    """将表格区域渲染为图片嵌入(暂不完美,后续可优化为 PPT Table 对象)"""
    bbox = table_info["bbox"]
    clip = fitz.Rect(bbox)

    zoom = dpi / 72
    mat = fitz.Matrix(zoom, zoom)
    pix = page.get_pixmap(matrix=mat, alpha=False, clip=clip)
    img_data = pix.tobytes("png")

    x0 = bbox[0] * scale_x
    y0 = bbox[1] * scale_y
    w = (bbox[2] - bbox[0]) * scale_x
    h = (bbox[3] - bbox[1]) * scale_y

    try:
        slide.shapes.add_picture(
            io.BytesIO(img_data),
            left=Emu(int(x0)),
            top=Emu(int(y0)),
            width=Emu(int(w)),
            height=Emu(int(h)),
        )
    except Exception:
        pass


def _font_available(font_name: str) -> bool:
    """简易判断字体是否在系统中可用

    实际应查询注册表或调用 Windows GDIEnumFontFamilies,
    但 fitz 在提取信息时用的是字体内部名,
    这里做一个简单关键词匹配。
    """
    # 已知常见字体名 -> Windows 系统字体文件
    SYSTEM_FONTS = {
        "宋体": "simsun.ttc",
        "黑体": "simhei.ttf",
        "微软雅黑": "msyh.ttc",
        "楷体": "simkai.ttf",
        "仿宋": "simfang.ttf",
        "Arial": "arial.ttf",
        "Times": "times.ttf",
        "Calibri": "calibri.ttf",
        "Tahoma": "tahoma.ttf",
        "Verdana": "verdana.ttf",
    }

    name_lower = font_name.lower()

    for key, file in SYSTEM_FONTS.items():
        if key.lower() in name_lower or name_lower in key.lower():
            return True

    # 检查 Windows 字体目录
    import glob
    windows_font_dir = os.path.join(os.environ.get("WINDIR", r"C:\Windows"), "Fonts")
    if os.path.isdir(windows_font_dir):
        try:
            for f in glob.glob(os.path.join(windows_font_dir, "*")):
                if font_name.lower() in os.path.basename(f).lower():
                    return True
        except Exception:
            pass

    return False


# ─── 入口 ──────────────────────────────────────

def main():
    try:
        # 1. 读取 stdin
        raw = sys.stdin.read()
        payload = json.loads(raw)
        inputs = payload["inputs"]
        options = payload.get("options", {})

        if not inputs:
            _error("没有输入文件")
            return

        pdf_path = inputs[0]

        if not os.path.isfile(pdf_path):
            _error(f"文件不存在: {pdf_path}")
            return

        # 2. 确定输出路径
        default_out_dir = os.path.dirname(pdf_path)
        output_dir = options.get("output_dir", default_out_dir)
        os.makedirs(output_dir, exist_ok=True)

        pdf_name = Path(pdf_path).stem
        output_path = os.path.join(output_dir, f"{pdf_name}.pptx")

        # 3. 执行转换
        _progress(5, "正在打开 PDF...")
        result = convert_pdf_to_pptx(pdf_path, output_path, options)

        # 4. 返回结果
        _result(
            output_files=[output_path],
            warnings=result.get("warnings", []),
            stats={"pageCount": result["page_count"]},
        )

    except json.JSONDecodeError as e:
        _error(f"JSON 解析错误: {e}")
    except Exception as e:
        import traceback
        tb = traceback.format_exc()
        _error(f"转换失败: {e}\n{tb}")


if __name__ == "__main__":
    main()
